"""
Generate AI Hackathon PPT — CS06 FHIR Converter
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_THEME_COLOR
import copy
from lxml import etree
import os

# ── Colors ──────────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x0A, 0x06, 0x2A)   # very dark navy/purple
MID_PURPLE   = RGBColor(0x2D, 0x1B, 0x69)   # mid purple
ACCENT_BLUE  = RGBColor(0x3B, 0x82, 0xF6)   # bright blue
ACCENT_PURP  = RGBColor(0x7C, 0x3A, 0xED)   # vivid purple
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xD6, 0xF1)
GOLD         = RGBColor(0xF5, 0xC5, 0x18)
CYAN         = RGBColor(0x06, 0xB6, 0xD4)
GREEN        = RGBColor(0x10, 0xB9, 0x81)
ORANGE       = RGBColor(0xF5, 0x9E, 0x0B)
FOOTER_GRAY  = RGBColor(0xA0, 0xAE, 0xC0)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def prs():
    p = Presentation()
    p.slide_width  = SLIDE_W
    p.slide_height = SLIDE_H
    return p


# ── XML helpers ─────────────────────────────────────────────────────────────
def add_grad_bg(slide, stops):
    """Add a linear gradient background fill to a slide via XML (solid fallback)."""
    # Use first stop color as solid background (most compatible approach)
    _, (r, g, b) = stops[0]
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

    # Also add a gradient rect shape that covers the whole slide
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    # Add a full-slide gradient rectangle shape
    from pptx.util import Emu as E
    grad_rect = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    grad_rect.line.fill.background()
    # Use XML to apply gradient fill to the shape
    sp = grad_rect._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        from lxml import etree as et2
        spPr = et2.SubElement(sp, qn('p:spPr'))
    # Remove existing fill
    for child in list(spPr):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('solidFill', 'gradFill', 'noFill', 'blipFill', 'pattFill'):
            spPr.remove(child)
    xml = f'<a:gradFill xmlns:a="{ns}"><a:gsLst>'
    for pos, (r2, g2, b2) in stops:
        xml += f'<a:gs pos="{pos}"><a:srgbClr val="{r2:02X}{g2:02X}{b2:02X}"/></a:gs>'
    xml += '</a:gsLst><a:lin ang="13500000" scaled="0"/></a:gradFill>'
    spPr.insert(0, etree.fromstring(xml))
    # Send to back
    slide.shapes._spTree.remove(grad_rect._element)
    slide.shapes._spTree.insert(2, grad_rect._element)


def set_cell_bg(cell, r, g, b):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    fill_xml = (
        f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:srgbClr val="{r:02X}{g:02X}{b:02X}"/>'
        f'</a:solidFill>'
    )
    existing = tcPr.find(qn('a:solidFill'))
    if existing is not None:
        tcPr.remove(existing)
    tcPr.insert(0, etree.fromstring(fill_xml))


def add_rect(slide, left, top, width, height, r, g, b, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.color.rgb = RGBColor(r, g, b)
    return shape


def add_textbox(slide, left, top, width, height, text, size, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.text_frame.word_wrap = wrap
    p = txb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def add_para(tf, text, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def footer(slide, right_text="2026 AI HACKATHON"):
    add_textbox(slide,
                Inches(0.3), Inches(7.1), Inches(8), Inches(0.35),
                "1455 LINCOLN PARKWAY EAST, 8TH FLOOR, ATLANTA, GA 30346  |  WWW.INNOVASOLUTIONS.COM",
                7, color=FOOTER_GRAY, align=PP_ALIGN.LEFT)
    add_textbox(slide,
                Inches(10.5), Inches(7.1), Inches(2.5), Inches(0.35),
                right_text, 7, color=FOOTER_GRAY, align=PP_ALIGN.RIGHT)
    # thin line above footer
    line = slide.shapes.add_connector(1, Inches(0.3), Inches(7.05),
                                       Inches(13.0), Inches(7.05))
    line.line.color.rgb = RGBColor(0x3B, 0x4A, 0x6B)
    line.line.width = Pt(0.5)


def innova_logo(slide, left=Inches(0.35), top=Inches(0.22)):
    """Simple text-based Innova logo placeholder."""
    txb = slide.shapes.add_textbox(left, top, Inches(2), Inches(0.5))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "innova"
    r1.font.size = Pt(16)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    r1.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = "®"
    r2.font.size = Pt(8)
    r2.font.color.rgb = WHITE
    r2.font.name = "Calibri"
    p2 = tf.add_paragraph()
    r3 = p2.add_run()
    r3.text = "SOLUTIONS"
    r3.font.size = Pt(8)
    r3.font.bold = True
    r3.font.color.rgb = LIGHT_GRAY
    r3.font.name = "Calibri"
    r3.font.letter_spacing = Pt(2)


# ── SLIDE 1: Title ───────────────────────────────────────────────────────────
def slide1(prs_obj):
    layout = prs_obj.slide_layouts[6]  # blank
    sl = prs_obj.slides.add_slide(layout)

    add_grad_bg(sl, [
        (0,      (0x08, 0x04, 0x22)),
        (35000,  (0x1A, 0x0A, 0x45)),
        (65000,  (0x2D, 0x1B, 0x69)),
        (100000, (0x0F, 0x0A, 0x35)),
    ])

    # Glow circles (decorative ovals)
    glow = sl.shapes.add_shape(9, Inches(7.5), Inches(1.5), Inches(4), Inches(4))  # oval=9
    glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x4C, 0x1D, 0x95)
    glow.line.fill.background()

    glow2 = sl.shapes.add_shape(9, Inches(1), Inches(4), Inches(3), Inches(3))
    glow2.fill.solid(); glow2.fill.fore_color.rgb = RGBColor(0x1E, 0x1B, 0x4B)
    glow2.line.fill.background()

    innova_logo(sl)

    # "AI HACKATHON" — big
    add_textbox(sl, Inches(1.5), Inches(1.8), Inches(10), Inches(1.6),
                "AI HACKATHON", 72, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    # Box: Case ID + Title
    box = add_rect(sl, Inches(2.3), Inches(3.3), Inches(8.7), Inches(1.05),
                   0x1E, 0x1B, 0x4B)
    box.line.color.rgb = ACCENT_BLUE
    box.line.width = Pt(1.2)
    txb = sl.shapes.add_textbox(Inches(2.4), Inches(3.35), Inches(8.5), Inches(1.0))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "2026 AI Hackathon  |  Case Study No. CS06"
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = '"FHIR Converter – Agentic Code Generation"'
    r2.font.size = Pt(18); r2.font.bold = True; r2.font.color.rgb = GOLD; r2.font.name = "Calibri"

    # Presenters box
    pbox = add_rect(sl, Inches(3.0), Inches(4.55), Inches(7.3), Inches(0.42),
                    0x0F, 0x0E, 0x2E)
    pbox.line.color.rgb = RGBColor(0x3B, 0x4A, 0x8B)
    pbox.line.width = Pt(0.8)
    add_textbox(sl, Inches(3.1), Inches(4.57), Inches(7.1), Inches(0.4),
                "Presenters – Sanjivi Syamsundar, Joshva Koilraj, Pavan Kandula, "
                "Gayathri Paladugu, Rishika Sambavi",
                10.5, italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    footer(sl)
    return sl


# ── SLIDE 2: Application Overview ────────────────────────────────────────────
def slide2(prs_obj):
    layout = prs_obj.slide_layouts[6]
    sl = prs_obj.slides.add_slide(layout)
    add_grad_bg(sl, [
        (0,      (0x08, 0x04, 0x22)),
        (50000,  (0x0D, 0x14, 0x32)),
        (100000, (0x0A, 0x06, 0x2A)),
    ])

    innova_logo(sl, Inches(11.8), Inches(0.18))

    # Header bar
    hdr = add_rect(sl, Inches(0), Inches(0.0), Inches(13.33), Inches(0.9),
                   0x1E, 0x1B, 0x4B)
    hdr.line.fill.background()
    add_textbox(sl, Inches(0.4), Inches(0.12), Inches(9), Inches(0.65),
                "Case Study No. CS06  |  Application Overview", 22,
                bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    add_textbox(sl, Inches(0.4), Inches(1.0), Inches(12.5), Inches(0.45),
                "HL7 & EHR → FHIR Converter  —  Healthcare Interoperability Tool",
                15, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    # What is it?
    add_textbox(sl, Inches(0.4), Inches(1.55), Inches(12.5), Inches(0.35),
                "What is it?", 13, bold=True, color=ACCENT_BLUE)
    add_textbox(sl, Inches(0.4), Inches(1.9), Inches(12.3), Inches(0.55),
                "A full-stack web application that converts healthcare messages between HL7 v2.x, "
                "EHR pipe-delimited records, and FHIR R4 — the modern healthcare interoperability standard.",
                11.5, color=LIGHT_GRAY, wrap=True)

    # 3 columns: key capabilities
    cols = [
        ("🔄  Conversion Modes",
         ["HL7 v2.x → FHIR R4",
          "FHIR R4 → HL7 v2.x",
          "Raw EHR → FHIR R4",
          "Supports ADT, ORM, ORU,\nSIU, MDM, DFT, VXU…"],
         ACCENT_BLUE),
        ("🤖  AI-Powered Engine",
         ["Claude Sonnet 4.6 (Anthropic)",
          "Groq (Llama 3.3 70B)",
          "Multi-step LLM extraction",
          "Handles any HL7 variant"],
         ACCENT_PURP),
        ("📄  Output & Reports",
         ["FHIR JSON & XML",
          "PDF Conversion Report",
          "CSV / Excel download",
          "Field-level mapping audit",
          "PHI masking / unmasking"],
         GREEN),
    ]
    col_x = [Inches(0.35), Inches(4.6), Inches(8.85)]
    for i, (title, items, clr) in enumerate(cols):
        cx = col_x[i]
        box = add_rect(sl, cx, Inches(2.6), Inches(4.0), Inches(3.95),
                       0x0F, 0x14, 0x2E)
        box.line.color.rgb = clr
        box.line.width = Pt(1.0)
        add_textbox(sl, cx + Inches(0.12), Inches(2.68), Inches(3.8), Inches(0.4),
                    title, 12, bold=True, color=clr)
        # bullet items
        txb = sl.shapes.add_textbox(cx + Inches(0.12), Inches(3.12),
                                     Inches(3.75), Inches(3.3))
        txb.text_frame.word_wrap = True
        first = True
        for item in items:
            if first:
                p = txb.text_frame.paragraphs[0]; first = False
            else:
                p = txb.text_frame.add_paragraph()
            p.space_before = Pt(4)
            r = p.add_run()
            r.text = "▸  " + item
            r.font.size = Pt(11); r.font.color.rgb = LIGHT_GRAY; r.font.name = "Calibri"

    footer(sl)
    return sl


# ── SLIDE 3: Conversion Details ───────────────────────────────────────────────
def slide3(prs_obj):
    layout = prs_obj.slide_layouts[6]
    sl = prs_obj.slides.add_slide(layout)
    add_grad_bg(sl, [
        (0,      (0x08, 0x04, 0x22)),
        (50000,  (0x0D, 0x14, 0x32)),
        (100000, (0x0A, 0x06, 0x2A)),
    ])
    innova_logo(sl, Inches(11.8), Inches(0.18))

    hdr = add_rect(sl, Inches(0), Inches(0.0), Inches(13.33), Inches(0.9),
                   0x1E, 0x1B, 0x4B)
    hdr.line.fill.background()
    add_textbox(sl, Inches(0.4), Inches(0.12), Inches(9), Inches(0.65),
                "Case Study No. CS06  |  Conversion Details & Report Downloads", 20,
                bold=True, color=WHITE)

    # ── Left column: supported segments ──
    add_textbox(sl, Inches(0.35), Inches(1.0), Inches(5.5), Inches(0.38),
                "Supported HL7 Segments → FHIR Resources", 12, bold=True, color=ACCENT_BLUE)

    seg_rows = [
        ("MSH", "MessageHeader / Organization"),
        ("PID", "Patient"),
        ("PV1", "Encounter  (class, location, attending)"),
        ("ORC", "ServiceRequest  (placer/filler/group IDs)"),
        ("OBR", "ServiceRequest + DiagnosticReport"),
        ("OBX", "Observation  (NM/CWE/ST value types)"),
        ("NK1", "RelatedPerson"),
        ("AL1", "AllergyIntolerance"),
        ("DG1", "Condition"),
        ("IN1", "Coverage"),
        ("RXA/RXR", "Immunization"),
        ("GT1", "RelatedPerson (guarantor)"),
        ("NTE", "Annotation note"),
    ]
    y = Inches(1.45)
    for seg, fhir in seg_rows:
        sb = add_rect(sl, Inches(0.35), y, Inches(1.1), Inches(0.27),
                      0x1E, 0x40, 0x7A)
        sb.line.fill.background()
        add_textbox(sl, Inches(0.37), y + Inches(0.02), Inches(1.05), Inches(0.25),
                    seg, 9, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
        add_textbox(sl, Inches(1.52), y, Inches(4.2), Inches(0.27),
                    fhir, 9.5, color=LIGHT_GRAY)
        y += Inches(0.31)

    # ── Right column: Report downloads ──
    add_textbox(sl, Inches(6.7), Inches(1.0), Inches(6.2), Inches(0.38),
                "Report & Download Options", 12, bold=True, color=GOLD)

    dl_items = [
        ("📄  FHIR JSON",
         "Full FHIR R4 Bundle — copy or download as .json\n"
         "Single patient:  Patient_SMITH_JOHN.json\n"
         "Multi-patient:  FHIR_Bundle_<msgType>.json"),
        ("📑  FHIR XML",
         "FHIR R4 Bundle in XML format (.xml)\n"
         "Fully schema-compliant HL7 FHIR R4 XML"),
        ("📊  PDF Report",
         "Branded conversion report with:\n"
         "  • Resource Summary table\n"
         "  • Per-resource field detail\n"
         "  • Field Mapping audit trail\n"
         "  Filename:  FHIR_Report_<msgType>_<date>.pdf"),
        ("📈  CSV / Excel",
         "Flat tabular view of all FHIR resources\n"
         "Filename:  FHIR_Export_<msgType>.csv / .xlsx\n"
         "One row per resource, columns per field"),
        ("🔒  PHI Masking",
         "Mask PID, SSN, NK1 before AI send\n"
         "🔓 Unmask restores original values\n"
         "Downloads respect mask/unmask state"),
    ]
    dy = Inches(1.45)
    for icon_title, desc in dl_items:
        bx = add_rect(sl, Inches(6.7), dy, Inches(6.25), Inches(0.82),
                      0x0F, 0x14, 0x2E)
        bx.line.color.rgb = RGBColor(0x2D, 0x3A, 0x6A)
        bx.line.width = Pt(0.7)
        add_textbox(sl, Inches(6.82), dy + Inches(0.03), Inches(6.0), Inches(0.28),
                    icon_title, 10.5, bold=True, color=WHITE)
        add_textbox(sl, Inches(6.82), dy + Inches(0.32), Inches(6.0), Inches(0.52),
                    desc, 9, color=LIGHT_GRAY, wrap=True)
        dy += Inches(0.9)

    footer(sl)
    return sl


# ── SLIDE 4: AI Conversion ────────────────────────────────────────────────────
def slide4(prs_obj):
    layout = prs_obj.slide_layouts[6]
    sl = prs_obj.slides.add_slide(layout)
    add_grad_bg(sl, [
        (0,      (0x08, 0x04, 0x22)),
        (50000,  (0x0D, 0x14, 0x32)),
        (100000, (0x0A, 0x06, 0x2A)),
    ])
    innova_logo(sl, Inches(11.8), Inches(0.18))

    hdr = add_rect(sl, Inches(0), Inches(0.0), Inches(13.33), Inches(0.9),
                   0x1E, 0x1B, 0x4B)
    hdr.line.fill.background()
    add_textbox(sl, Inches(0.4), Inches(0.12), Inches(10), Inches(0.65),
                "Case Study No. CS06  |  AI-Powered Conversion", 20,
                bold=True, color=WHITE)

    # Intro
    add_textbox(sl, Inches(0.4), Inches(1.0), Inches(12.5), Inches(0.38),
                "Dual AI Provider Architecture — Claude Sonnet 4.6 (Anthropic) & Groq (Llama 3.3 70B)",
                13, bold=True, color=GOLD)
    add_textbox(sl, Inches(0.4), Inches(1.42), Inches(12.5), Inches(0.35),
                "Toggle AI Mode ON to unlock LLM-powered conversion for complex, real-world HL7 and EHR messages "
                "that rule-based parsers cannot handle.",
                11, color=LIGHT_GRAY, wrap=True)

    # ── How AI works: pipeline ──
    add_textbox(sl, Inches(0.4), Inches(1.88), Inches(6.0), Inches(0.32),
                "4-Step Extraction Pipeline", 11.5, bold=True, color=ACCENT_BLUE)

    steps = [
        ("Step 1", "Demographics", "PID→Patient, PV1→Encounter,\nNK1→RelatedPerson", ACCENT_BLUE),
        ("Step 2", "Clinical",     "ORC→ServiceRequest, OBR→DiagnosticReport,\nOBX→Observation, RXA→Immunization", ACCENT_PURP),
        ("Step 3", "Admin",        "AL1→AllergyIntolerance, DG1→Condition,\nIN1→Coverage, Z-segments→Parameters", GREEN),
        ("Step 4", "Field Maps",   "Full field-level HL7→FHIR mapping\naudit trail returned to UI", ORANGE),
    ]
    sx = Inches(0.35)
    for i, (num, title, desc, clr) in enumerate(steps):
        bx = add_rect(sl, sx, Inches(2.28), Inches(3.0), Inches(1.5),
                      0x0F, 0x14, 0x2E)
        bx.line.color.rgb = clr; bx.line.width = Pt(1.2)
        nb = add_rect(sl, sx, Inches(2.28), Inches(0.7), Inches(0.36),
                      *[c for c in clr])
        nb.line.fill.background()
        add_textbox(sl, sx + Inches(0.03), Inches(2.29), Inches(0.65), Inches(0.32),
                    num, 8, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_textbox(sl, sx, Inches(2.66), Inches(3.0), Inches(0.32),
                    title, 11, bold=True, color=clr, align=PP_ALIGN.CENTER)
        add_textbox(sl, sx + Inches(0.08), Inches(3.0), Inches(2.85), Inches(0.75),
                    desc, 9.5, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, wrap=True)
        if i < 3:
            arr = sl.shapes.add_connector(1,
                sx + Inches(3.0), Inches(3.03),
                sx + Inches(3.35), Inches(3.03))
            arr.line.color.rgb = RGBColor(0x4B, 0x5A, 0x8B)
            arr.line.width = Pt(1.5)
        sx += Inches(3.2)

    # ── Right panel: AI benefits ──
    add_textbox(sl, Inches(0.4), Inches(3.92), Inches(5.5), Inches(0.32),
                "Why AI for Healthcare Conversion?", 11.5, bold=True, color=ACCENT_BLUE)

    benefits = [
        ("🧠  Handles any HL7 variant",
         "Rule-based parsers fail on non-standard or custom Z-segments. "
         "LLMs understand clinical context and infer correct mappings."),
        ("⚡  Multi-provider flexibility",
         "Groq (free tier, fast) for development; Claude Sonnet 4.6 for "
         "production accuracy. Radio-button switch in UI — no code change."),
        ("🔒  PHI-safe by design",
         "Mask PHI checkbox replaces patient name, DOB, MRN, SSN, NK1 "
         "data before the message leaves your browser. Unmask locally."),
        ("🛠️  JSON repair built-in",
         "json-repair library recovers malformed LLM output. "
         "_unwrap_to_list() handles Bundle vs array responses from any LLM."),
    ]
    by = Inches(4.28)
    for icon_title, desc in benefits:
        bx2 = add_rect(sl, Inches(0.35), by, Inches(6.1), Inches(0.65),
                       0x0F, 0x14, 0x2E)
        bx2.line.color.rgb = RGBColor(0x2D, 0x3A, 0x6A); bx2.line.width = Pt(0.6)
        add_textbox(sl, Inches(0.48), by + Inches(0.03), Inches(5.8), Inches(0.25),
                    icon_title, 10, bold=True, color=WHITE)
        add_textbox(sl, Inches(0.48), by + Inches(0.28), Inches(5.8), Inches(0.35),
                    desc, 9, color=LIGHT_GRAY, wrap=True)
        by += Inches(0.73)

    # ── Right panel: tech stack ──
    add_textbox(sl, Inches(6.9), Inches(3.92), Inches(6.0), Inches(0.32),
                "Tech Stack", 11.5, bold=True, color=GOLD)
    stack = [
        ("Backend",  "Python · FastAPI · Uvicorn", ACCENT_BLUE),
        ("AI",       "Anthropic SDK (Claude) · Groq SDK · json-repair", ACCENT_PURP),
        ("Frontend", "Vanilla JS · CSS3 · HTML5 · jsPDF · SheetJS", GREEN),
        ("HL7",      "python-hl7 · Custom ORM/ORU/ADT/SIU converters", CYAN),
        ("Deploy",   "localhost:8000  (Windows/Linux/Mac)", ORANGE),
    ]
    ty = Inches(4.28)
    for label, val, clr in stack:
        lb = add_rect(sl, Inches(6.9), ty, Inches(1.2), Inches(0.36),
                      0x1E, 0x28, 0x50)
        lb.line.fill.background()
        add_textbox(sl, Inches(6.92), ty + Inches(0.04), Inches(1.15), Inches(0.28),
                    label, 9, bold=True, color=clr, align=PP_ALIGN.CENTER)
        add_textbox(sl, Inches(8.18), ty + Inches(0.04), Inches(4.7), Inches(0.28),
                    val, 10, color=LIGHT_GRAY)
        ty += Inches(0.58)

    footer(sl)
    return sl


# ── SLIDE 5: Demo Video ───────────────────────────────────────────────────────
def slide5(prs_obj):
    layout = prs_obj.slide_layouts[6]
    sl = prs_obj.slides.add_slide(layout)
    add_grad_bg(sl, [
        (0,      (0x08, 0x04, 0x22)),
        (50000,  (0x0D, 0x14, 0x32)),
        (100000, (0x0A, 0x06, 0x2A)),
    ])
    innova_logo(sl, Inches(11.8), Inches(0.18))

    hdr = add_rect(sl, Inches(0), Inches(0.0), Inches(13.33), Inches(0.9),
                   0x1E, 0x1B, 0x4B)
    hdr.line.fill.background()
    add_textbox(sl, Inches(0.4), Inches(0.12), Inches(10), Inches(0.65),
                "Case Study No. CS06  |  Live Demo", 20, bold=True, color=WHITE)

    # Video placeholder box
    vbox = add_rect(sl, Inches(1.2), Inches(1.05), Inches(10.93), Inches(5.3),
                    0x0F, 0x14, 0x2E)
    vbox.line.color.rgb = ACCENT_BLUE; vbox.line.width = Pt(2)

    # Play button circle
    play = sl.shapes.add_shape(9, Inches(5.9), Inches(2.8), Inches(1.5), Inches(1.5))
    play.fill.solid(); play.fill.fore_color.rgb = ACCENT_BLUE
    play.line.fill.background()

    add_textbox(sl, Inches(6.1), Inches(2.98), Inches(0.9), Inches(0.6),
                "▶", 32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(sl, Inches(1.2), Inches(4.45), Inches(10.93), Inches(0.45),
                "▶  INSERT VIDEO HERE  —  Demo_FHIR_Converter_CS06.mp4  (~2 min)",
                14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    add_textbox(sl, Inches(1.5), Inches(4.95), Inches(10.33), Inches(0.9),
                "Recording steps:  Open localhost:8000  →  Paste ORM HL7 sample  →  Convert (manual)  →"
                "  Show FHIR JSON / PDF Report  →  Enable AI Mode (Claude)  →  Convert EHR sample  →"
                "  Download Excel  →  Demonstrate PHI Mask / Unmask",
                9.5, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, wrap=True)

    # How to embed
    hint = add_rect(sl, Inches(1.2), Inches(5.92), Inches(10.93), Inches(0.45),
                    0x1A, 0x24, 0x4A)
    hint.line.fill.background()
    add_textbox(sl, Inches(1.35), Inches(5.95), Inches(10.6), Inches(0.36),
                "💡  To embed:  Insert → Video → This Device → select Demo_FHIR_Converter_CS06.mp4",
                10, color=CYAN, align=PP_ALIGN.CENTER)

    footer(sl)
    return sl


# ── Build & save ─────────────────────────────────────────────────────────────
if __name__ == "__main__":
    p = prs()
    slide1(p)
    slide2(p)
    slide3(p)
    slide4(p)
    slide5(p)

    out = os.path.join(
        r"c:\Users\jkoilraj\OneDrive - InnovaSolutions\Documents\AI\HL7toFHIR",
        "CS06_FHIR_Converter_Hackathon.pptx"
    )
    p.save(out)
    print(f"Saved: {out}")
